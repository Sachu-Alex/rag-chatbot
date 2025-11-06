"""
Advanced Document Processing System

This module provides comprehensive document processing capabilities for the RAG system.
It handles multiple file formats (PDF, DOCX, PPTX, TXT, Markdown) and converts them
into structured chunks suitable for vector storage and retrieval.

Key Features:
- Multi-format document support (PDF, DOCX, PPTX, TXT, MD)
- Intelligent text extraction with metadata preservation
- Table extraction from documents
- Configurable text chunking with overlap
- File caching for performance optimization
- Comprehensive error handling and logging

Classes:
    DocumentProcessor: Main class for processing documents into chunks

Usage Example:
    processor = DocumentProcessor()
    documents = processor.process_file("document.pdf")
    for doc in documents:
        # Process each document chunk
        pass
"""
import os
import logging
from typing import List, Dict, Any, Optional, Union
from pathlib import Path
import hashlib

# Document processing imports
import PyPDF2
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd

# LangChain imports
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter

# Configuration
from config import config

logging.basicConfig(level=getattr(logging, config.log_level))
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """
    Advanced Document Processor with Multi-Format Support
    
    This class handles the complete document processing pipeline from file ingestion
    to structured text chunks ready for vector storage. It supports multiple file
    formats and provides configurable text splitting with intelligent chunking.
    
    Key Capabilities:
    - PDF processing with table extraction using pdfplumber
    - DOCX document parsing with table support
    - PowerPoint (PPTX) slide text extraction
    - Plain text and Markdown file processing
    - File hash-based caching for performance
    - Configurable text chunking with overlap
    - Comprehensive metadata extraction
    
    Attributes:
        config: Document processing configuration
        text_splitter: LangChain text splitter for creating chunks
        document_cache: Cache for processed documents (hash-based)
    """
    
    def __init__(self, processing_config=None):
        """
        Initialize the document processor with configuration.
        
        Args:
            processing_config: Optional configuration object. If None, uses default config.
                             Should contain chunk_size, chunk_overlap, separators, etc.
        """
        self.config = processing_config or config.document_processing
        
        # Initialize text splitter
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.config.chunk_size,
            chunk_overlap=self.config.chunk_overlap,
            separators=self.config.separators,
            length_function=len
        )
        
        # Document metadata cache
        self.document_cache = {}
    
    def compute_file_hash(self, file_path: str) -> str:
        """
        Compute SHA-256 hash of a file for caching purposes.
        
        This method creates a unique fingerprint for each file, allowing the system
        to cache processed results and avoid reprocessing unchanged files.
        
        Args:
            file_path (str): Path to the file to hash
            
        Returns:
            str: SHA-256 hexadecimal hash of the file content
            
        Note:
            Returns empty string if file cannot be read or hashed
        """
        hash_sha256 = hashlib.sha256()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_sha256.update(chunk)
            return hash_sha256.hexdigest()
        except Exception as e:
            logger.error(f"Error computing hash for {file_path}: {e}")
            return ""
    
    def extract_text_from_pdf(self, file_path: str) -> tuple[str, Dict[str, Any]]:
        """
        Extract text content from PDF files with advanced features.
        
        This method uses pdfplumber as the primary extraction tool for better
        text and table extraction, with PyPDF2 as a fallback. It extracts:
        - Main text content from all pages
        - Table data (if extract_tables is enabled)
        - Document metadata (page count, images, tables)
        
        Args:
            file_path (str): Path to the PDF file
            
        Returns:
            tuple[str, Dict[str, Any]]: 
                - Extracted text content with page markers
                - Metadata dictionary containing extraction info
                
        Raises:
            Exception: If both pdfplumber and PyPDF2 extraction fail
        """
        try:
            text = ""
            metadata = {
                "pages": 0,
                "has_images": False,
                "has_tables": False,
                "extraction_method": "pdfplumber"
            }
            
            # Use pdfplumber for better text extraction
            with pdfplumber.open(file_path) as pdf:
                metadata["pages"] = len(pdf.pages)
                
                for i, page in enumerate(pdf.pages):
                    # Extract text
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n--- Page {i+1} ---\n{page_text}\n"
                    
                    # Check for images and tables
                    if not metadata["has_images"] and page.images:
                        metadata["has_images"] = True
                    
                    if not metadata["has_tables"] and page.extract_tables():
                        metadata["has_tables"] = True
                        
                        # Extract table data if enabled
                        if self.config.extract_tables:
                            tables = page.extract_tables()
                            for j, table in enumerate(tables):
                                if table:
                                    text += f"\n--- Table {j+1} on Page {i+1} ---\n"
                                    # Convert table to text format
                                    for row in table:
                                        if row:
                                            text += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
                                    text += "\n"
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"Error extracting from PDF {file_path}: {e}")
            # Fallback to PyPDF2
            try:
                text = ""
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                
                metadata = {
                    "pages": len(pdf_reader.pages),
                    "extraction_method": "pypdf2_fallback"
                }
                return text, metadata
                
            except Exception as e2:
                logger.error(f"Fallback extraction also failed: {e2}")
                raise
    
    def extract_text_from_docx(self, file_path: str) -> tuple[str, Dict[str, Any]]:
        """Extract text from DOCX with metadata."""
        try:
            doc = DocxDocument(file_path)
            
            # Extract paragraphs
            text = ""
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            # Extract tables if enabled
            if self.config.extract_tables and doc.tables:
                text += "\n--- Tables ---\n"
                for i, table in enumerate(doc.tables):
                    text += f"\nTable {i+1}:\n"
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        text += row_text + "\n"
                    text += "\n"
            
            metadata = {
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables),
                "extraction_method": "python-docx"
            }
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"Error extracting from DOCX {file_path}: {e}")
            raise
    
    def extract_text_from_pptx(self, file_path: str) -> tuple[str, Dict[str, Any]]:
        """Extract text from PowerPoint with metadata."""
        try:
            prs = Presentation(file_path)
            text = ""
            slide_count = 0
            
            for i, slide in enumerate(prs.slides):
                slide_text = f"\n--- Slide {i+1} ---\n"
                slide_content = ""
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content += shape.text + "\n"
                
                if slide_content.strip():
                    text += slide_text + slide_content
                    slide_count += 1
            
            metadata = {
                "total_slides": len(prs.slides),
                "slides_with_text": slide_count,
                "extraction_method": "python-pptx"
            }
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"Error extracting from PPTX {file_path}: {e}")
            raise
    
    def extract_text_from_txt(self, file_path: str) -> tuple[str, Dict[str, Any]]:
        """Extract text from plain text file."""
        try:
            encodings = ['utf-8', 'latin-1', 'ascii', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text = file.read()
                    
                    metadata = {
                        "encoding": encoding,
                        "lines": len(text.splitlines()),
                        "extraction_method": "direct"
                    }
                    
                    return text, metadata
                    
                except UnicodeDecodeError:
                    continue
            
            raise ValueError(f"Could not decode file {file_path} with any supported encoding")
            
        except Exception as e:
            logger.error(f"Error extracting from TXT {file_path}: {e}")
            raise
    
    def extract_text_from_markdown(self, file_path: str) -> tuple[str, Dict[str, Any]]:
        """Extract text from Markdown file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            
            # Count markdown elements
            lines = text.splitlines()
            headers = len([line for line in lines if line.startswith('#')])
            code_blocks = text.count('```')
            
            metadata = {
                "lines": len(lines),
                "headers": headers,
                "code_blocks": code_blocks // 2,  # Pairs of ```
                "extraction_method": "direct"
            }
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"Error extracting from Markdown {file_path}: {e}")
            raise
    
    def process_file(self, file_path: str) -> List[Document]:
        """Process a single file and return chunked documents."""
        file_path = str(Path(file_path).resolve())
        file_extension = Path(file_path).suffix.lower()
        
        logger.info(f"Processing file: {file_path}")
        
        # Check if file is supported
        if file_extension not in self.config.supported_extensions:
            raise ValueError(f"Unsupported file type: {file_extension}")
        
        # Compute file hash for caching
        file_hash = self.compute_file_hash(file_path)
        cache_key = f"{file_path}:{file_hash}"
        
        if cache_key in self.document_cache:
            logger.info(f"Using cached result for {file_path}")
            return self.document_cache[cache_key]
        
        try:
            # Extract text based on file type
            if file_extension == '.pdf':
                text, extraction_metadata = self.extract_text_from_pdf(file_path)
            elif file_extension == '.docx':
                text, extraction_metadata = self.extract_text_from_docx(file_path)
            elif file_extension == '.pptx':
                text, extraction_metadata = self.extract_text_from_pptx(file_path)
            elif file_extension == '.txt':
                text, extraction_metadata = self.extract_text_from_txt(file_path)
            elif file_extension == '.md':
                text, extraction_metadata = self.extract_text_from_markdown(file_path)
            else:
                raise ValueError(f"Handler not implemented for {file_extension}")
            
            if not text.strip():
                logger.warning(f"No text extracted from {file_path}")
                return []
            
            # Split text into chunks
            text_chunks = self.text_splitter.split_text(text)
            
            # Create documents with metadata
            documents = []
            for i, chunk in enumerate(text_chunks):
                doc_metadata = {
                    "source": file_path,
                    "filename": Path(file_path).name,
                    "file_type": file_extension,
                    "chunk_id": i,
                    "total_chunks": len(text_chunks),
                    "file_hash": file_hash,
                    **extraction_metadata
                }
                
                doc = Document(
                    page_content=chunk,
                    metadata=doc_metadata
                )
                documents.append(doc)
            
            # Cache results
            self.document_cache[cache_key] = documents
            
            logger.info(f"Successfully processed {file_path}: {len(documents)} chunks created")
            return documents
            
        except Exception as e:
            logger.error(f"Failed to process {file_path}: {e}")
            raise
    
    def process_directory(self, directory_path: str) -> List[Document]:
        """Process all supported files in a directory."""
        directory = Path(directory_path)
        if not directory.exists():
            raise ValueError(f"Directory does not exist: {directory_path}")
        
        all_documents = []
        processed_files = 0
        
        # Find all supported files
        for extension in self.config.supported_extensions:
            pattern = f"**/*{extension}"
            files = list(directory.glob(pattern))
            
            for file_path in files:
                try:
                    documents = self.process_file(str(file_path))
                    all_documents.extend(documents)
                    processed_files += 1
                except Exception as e:
                    logger.error(f"Skipping {file_path}: {e}")
                    continue
        
        logger.info(f"Processed {processed_files} files, created {len(all_documents)} document chunks")
        return all_documents
    
    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """Get information about a file without processing it."""
        file_path = Path(file_path)
        
        info = {
            "filename": file_path.name,
            "file_type": file_path.suffix.lower(),
            "file_size": file_path.stat().st_size if file_path.exists() else 0,
            "supported": file_path.suffix.lower() in self.config.supported_extensions,
            "exists": file_path.exists()
        }
        
        return info
    
    def clear_cache(self):
        """Clear the document processing cache."""
        self.document_cache.clear()
        logger.info("Document processing cache cleared")
    
    def get_cache_stats(self) -> Dict[str, Any]:
        """Get statistics about the document cache."""
        return {
            "cached_documents": len(self.document_cache),
            "total_chunks": sum(len(docs) for docs in self.document_cache.values())
        }